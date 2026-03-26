"""
Document Ingestion Service.
Handles file upload, GridFS storage, and text extraction from PDF/DOCX/PPTX/images.

How to run locally:
    - Install dependencies: pip install PyMuPDF python-docx python-pptx pytesseract Pillow
    - For OCR: install Tesseract from https://github.com/tesseract-ocr/tesseract
    - Set TESSERACT_CMD env var if tesseract is not on PATH

How to switch from mocked LLM to real model:
    - This module does NOT use an LLM — it only extracts text.
"""

import os
import io
import re
import logging
from typing import Tuple, List, Dict, Optional
from datetime import datetime

import fitz  # PyMuPDF
from bson import ObjectId
from motor.motor_asyncio import AsyncIOMotorDatabase, AsyncIOMotorGridFSBucket

logger = logging.getLogger(__name__)


# ── Text Extraction ──────────────────────────────────────────────────────────

def extract_text_from_pdf(file_bytes: bytes) -> Tuple[str, int]:
    """
    Extract text from PDF using PyMuPDF.
    Returns (extracted_text, page_count).
    Handles tables, headers, and multi-column layouts.
    """
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages_text = []
    for page in doc:
        # Use "text" extraction with sorting for multi-column support
        text = page.get_text("text", sort=True)
        if text.strip():
            pages_text.append(text)
        else:
            # Fallback: try OCR if page is image-based
            # (requires pytesseract — handled in extract_text_from_image)
            pix = page.get_pixmap(dpi=200)
            img_bytes = pix.tobytes("png")
            ocr_text = _ocr_image_bytes(img_bytes)
            if ocr_text.strip():
                pages_text.append(ocr_text)

    page_count = len(doc)
    doc.close()
    full_text = "\n\n".join(pages_text)
    return _clean_text(full_text), page_count


def extract_text_from_docx(file_bytes: bytes) -> Tuple[str, int]:
    """
    Extract text from DOCX using python-docx.
    Returns (extracted_text, estimated_page_count).
    """
    try:
        from docx import Document
    except ImportError:
        raise ImportError("python-docx is required: pip install python-docx")

    doc = Document(io.BytesIO(file_bytes))
    paragraphs = []
    for para in doc.paragraphs:
        if para.text.strip():
            # Preserve heading structure as section markers
            if para.style and para.style.name.startswith("Heading"):
                paragraphs.append(f"\n## {para.text.strip()}\n")
            else:
                paragraphs.append(para.text.strip())

    # Also extract from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)

    full_text = "\n".join(paragraphs)
    # Estimate page count: ~3000 chars per page
    page_count = max(1, len(full_text) // 3000)
    return _clean_text(full_text), page_count


def extract_text_from_pptx(file_bytes: bytes) -> Tuple[str, int]:
    """
    Extract text from PPTX using python-pptx.
    Returns (extracted_text, slide_count).
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx is required: pip install python-pptx")

    prs = Presentation(io.BytesIO(file_bytes))
    slides_text = []
    for i, slide in enumerate(prs.slides):
        slide_parts = [f"\n## Slide {i+1}\n"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_parts.append(text)
            # Extract from tables in slides
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        slide_parts.append(row_text)
        slides_text.append("\n".join(slide_parts))

    full_text = "\n\n".join(slides_text)
    return _clean_text(full_text), len(prs.slides)


def extract_text_from_image(file_bytes: bytes) -> Tuple[str, int]:
    """
    Extract text from image using Tesseract OCR.
    Returns (extracted_text, 1).
    """
    text = _ocr_image_bytes(file_bytes)
    return _clean_text(text), 1


def _ocr_image_bytes(image_bytes: bytes) -> str:
    """Run Tesseract OCR on image bytes."""
    try:
        import pytesseract
        from PIL import Image

        # Configure tesseract path if set via env
        tesseract_cmd = os.getenv("TESSERACT_CMD")
        if tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = tesseract_cmd

        image = Image.open(io.BytesIO(image_bytes))
        text = pytesseract.image_to_string(image, lang="eng+fra")
        return text
    except ImportError:
        logger.warning("pytesseract not installed. OCR disabled. pip install pytesseract Pillow")
        return ""
    except Exception as e:
        logger.warning(f"OCR failed: {e}")
        return ""


def _clean_text(text: str) -> str:
    """Clean extracted text: normalize whitespace, remove artifacts."""
    # Normalize line endings
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    # Remove excessive blank lines (3+ → 2)
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Remove page numbers (standalone numbers on a line)
    text = re.sub(r"^\s*\d{1,4}\s*$", "", text, flags=re.MULTILINE)
    # Remove header/footer artifacts (repeated short lines)
    text = re.sub(r"(?m)^.{1,15}$\n(?=.{1,15}$)", "", text)
    # Normalize whitespace within lines
    text = re.sub(r"[ \t]+", " ", text)
    return text.strip()


def detect_sections(text: str) -> List[Dict[str, str]]:
    """
    Detect section headings in extracted text.
    Looks for markdown-style headings (##) or ALL-CAPS lines.
    Returns list of {"title": str, "start_pos": int}.
    """
    sections = []
    lines = text.split("\n")
    pos = 0
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("## "):
            sections.append({"title": stripped[3:].strip(), "start_pos": pos})
        elif (
            len(stripped) > 3
            and len(stripped) < 100
            and stripped.isupper()
            and not stripped.isdigit()
        ):
            sections.append({"title": stripped.title(), "start_pos": pos})
        pos += len(line) + 1

    if not sections:
        sections.append({"title": "Main Content", "start_pos": 0})

    return sections


def get_file_type(filename: str) -> str:
    """Determine file type from filename extension."""
    ext = os.path.splitext(filename)[1].lower()
    type_map = {
        ".pdf": "pdf",
        ".docx": "docx",
        ".doc": "docx",
        ".pptx": "pptx",
        ".ppt": "pptx",
        ".png": "image",
        ".jpg": "image",
        ".jpeg": "image",
        ".tiff": "image",
        ".bmp": "image",
    }
    return type_map.get(ext, "pdf")


# ── Ingestion Pipeline ───────────────────────────────────────────────────────

async def store_file_in_gridfs(
    db: AsyncIOMotorDatabase,
    file_bytes: bytes,
    filename: str,
    content_type: str = "application/octet-stream"
) -> str:
    """
    Store raw file in MongoDB GridFS.
    Returns the GridFS file_id as string.
    """
    fs = AsyncIOMotorGridFSBucket(db)
    file_id = await fs.upload_from_stream(
        filename,
        io.BytesIO(file_bytes),
        metadata={"content_type": content_type, "uploaded_at": datetime.utcnow().isoformat()}
    )
    logger.info(f"Stored file '{filename}' in GridFS: {file_id}")
    return str(file_id)


async def ingest_document(
    db: AsyncIOMotorDatabase,
    file_bytes: bytes,
    filename: str,
    uploaded_by: str = "system"
) -> Dict:
    """
    Full ingestion pipeline:
    1. Detect file type
    2. Store raw file in GridFS
    3. Extract text
    4. Detect sections
    5. Create quiz_documents record

    Returns the created document dict with _id.
    """
    file_type = get_file_type(filename)

    # 1. Store in GridFS
    gridfs_id = await store_file_in_gridfs(db, file_bytes, filename)

    # 2. Extract text
    extractors = {
        "pdf": extract_text_from_pdf,
        "docx": extract_text_from_docx,
        "pptx": extract_text_from_pptx,
        "image": extract_text_from_image,
    }
    extractor = extractors.get(file_type, extract_text_from_pdf)

    try:
        text, page_count = extractor(file_bytes)
    except Exception as e:
        logger.error(f"Text extraction failed for '{filename}': {e}")
        # Create document record with error status
        doc = {
            "title": os.path.splitext(filename)[0],
            "filename": filename,
            "file_type": file_type,
            "gridfs_file_id": gridfs_id,
            "uploaded_by": uploaded_by,
            "uploaded_at": datetime.utcnow(),
            "status": "error",
            "total_chunks": 0,
            "total_tokens": 0,
            "sections": [],
            "metadata": {"page_count": 0, "language": "en", "category": None},
            "extracted_text": "",
        }
        result = await db.quiz_documents.insert_one(doc)
        doc["_id"] = result.inserted_id
        return doc

    # 3. Detect sections
    sections = detect_sections(text)

    # 4. Create document record
    doc = {
        "title": os.path.splitext(filename)[0],
        "filename": filename,
        "file_type": file_type,
        "gridfs_file_id": gridfs_id,
        "uploaded_by": uploaded_by,
        "uploaded_at": datetime.utcnow(),
        "status": "processing",
        "total_chunks": 0,
        "total_tokens": 0,
        "sections": [{"title": s["title"], "start_chunk": 0, "end_chunk": 0} for s in sections],
        "metadata": {"page_count": page_count, "language": "en", "category": None},
        "extracted_text": text,  # Stored temporarily for chunking, can be removed after
    }

    result = await db.quiz_documents.insert_one(doc)
    doc["_id"] = result.inserted_id
    logger.info(f"Ingested document '{filename}' → {result.inserted_id} ({len(text)} chars, {page_count} pages)")
    return doc
